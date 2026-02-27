#!/usr/bin/env python3
"""Build Cut the Crap — General Audience PowerPoint deck (Feb 2026 edition)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT    = RGBColor(0x00, 0x96, 0xD6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY  = RGBColor(0x33, 0x33, 0x33)
TABLE_HDR = RGBColor(0x0A, 0x1E, 0x3D)
TABLE_ALT = RGBColor(0xE8, 0xF4, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────────────────────

def _add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def _set(tf, text, size=18, bold=False, color=MED_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return p

def _add_para(tf, text, size=18, bold=False, color=MED_GRAY, alignment=PP_ALIGN.LEFT, space_before=Pt(6), bullet=False, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if bullet:
        p.level = 0
    return p

def add_section_slide(title, subtitle, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide, DARK_BG)
    # accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(0.15), Inches(1.6))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
    # title
    tb = _tb(slide, 1.3, 2.6, 10, 1.2)
    _set(tb.text_frame, title, size=40, bold=True, color=WHITE)
    # subtitle
    if subtitle:
        tb2 = _tb(slide, 1.3, 3.9, 10, 0.8)
        _set(tb2.text_frame, subtitle, size=22, color=ACCENT)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_content_slide(title, bullets, notes="", two_col=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    # title bar
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.1))
    rect.fill.solid(); rect.fill.fore_color.rgb = DARK_BG; rect.line.fill.background()
    _set(rect.text_frame, "   " + title, size=28, bold=True, color=WHITE)
    rect.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    rect.text_frame.paragraphs[0].space_before = Pt(14)

    if not two_col:
        tb = _tb(slide, 0.8, 1.4, 11.5, 5.5)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                _set(tf, "•  " + b, size=20, color=MED_GRAY)
            else:
                _add_para(tf, "•  " + b, size=20, color=MED_GRAY, space_before=Pt(12))
    else:
        mid = len(bullets) // 2
        left_b = bullets[:mid]
        right_b = bullets[mid:]
        tb = _tb(slide, 0.8, 1.4, 5.5, 5.5)
        tf = tb.text_frame; tf.word_wrap = True
        for i, b in enumerate(left_b):
            if i == 0:
                _set(tf, "•  " + b, size=20, color=MED_GRAY)
            else:
                _add_para(tf, "•  " + b, size=20, color=MED_GRAY, space_before=Pt(12))
        tb2 = _tb(slide, 6.8, 1.4, 5.5, 5.5)
        tf2 = tb2.text_frame; tf2.word_wrap = True
        for i, b in enumerate(right_b):
            if i == 0:
                _set(tf2, "•  " + b, size=20, color=MED_GRAY)
            else:
                _add_para(tf2, "•  " + b, size=20, color=MED_GRAY, space_before=Pt(12))
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_table_slide(title, headers, rows, notes="", col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    # title bar
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.1))
    rect.fill.solid(); rect.fill.fore_color.rgb = DARK_BG; rect.line.fill.background()
    _set(rect.text_frame, "   " + title, size=28, bold=True, color=WHITE)
    rect.text_frame.paragraphs[0].space_before = Pt(14)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = Inches(11.5)
    tbl_h = Inches(min(5.0, 0.5 + 0.45 * n_rows))
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.9), Inches(1.4), tbl_w, tbl_h)
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    def _cell(r, c, text, is_header=False):
        cell = tbl.cell(r, c)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "Calibri"
            if is_header:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = MED_GRAY
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HDR
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else TABLE_ALT

    for c, h in enumerate(headers):
        _cell(0, c, h, is_header=True)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            _cell(r + 1, c, val)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

# ── 1. TITLE ─────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, DARK_BG)
tb = _tb(slide, 1.0, 1.8, 11, 1.5)
_set(tb.text_frame, "Cut the Crap", size=54, bold=True, color=WHITE)
tb2 = _tb(slide, 1.0, 3.2, 11, 1.0)
_set(tb2.text_frame, "Everything You Need to Know About AI", size=30, color=ACCENT)
tb3 = _tb(slide, 1.0, 4.5, 11, 0.6)
_set(tb3.text_frame, "General Audience  •  4-Hour Workshop  •  February 2026", size=18, color=WHITE)
# accent bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.2), Inches(3), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

# ── 2. Agenda ────────────────────────────────────────────────────────────────
add_content_slide("Today's Agenda", [
    "Part 1 — The Landscape: ChatGPT vs Claude vs Gemini (30 min)",
    "Part 2 — How It Works: Models, pricing, API keys (30 min)",
    "Part 3 — What AI Can Do: Images, voice, vision demos (20 min)",
    "Part 4 — What AI Gets Wrong: Hallucinations & privacy (20 min)",
    "Part 5 — Hands-On: Try all three AIs yourself (40 min)",
    "Part 6 — OpenClaw Setup: Markdown, Linux basics, install (50 min)",
])

# ══ PART 1: THE LANDSCAPE ════════════════════════════════════════════════════

add_section_slide("Part 1: The Landscape", "There's more than ChatGPT — meet the Big Three")

# Slide: Three AIs
add_content_slide("There's More Than ChatGPT", [
    "ChatGPT (OpenAI) — The one everyone knows. Best ecosystem, image gen built-in.",
    "Claude (Anthropic) — Careful, great at long docs & writing. Less hallucination.",
    "Gemini (Google) — Integrated with Google services, current info, good at research.",
    "None is \"the best\" — they're like different restaurants for different cravings.",
    "By end of today you'll know when to use which one.",
],
notes="Raise your hand if you've used ChatGPT. Now keep up if you've also tried Claude or Gemini. Most hands drop — that's normal. ChatGPT had a massive head start (Nov 2022). But there are now three serious players, each good at different things."
)

# Slide: Live Comparison
add_content_slide("Live Comparison — Same Prompt, 3 AIs", [
    "Demo 1 (Creative): \"Write a 4-sentence horror story about a smart home\"",
    "Demo 2 (Factual): \"Three biggest Canadian tech acquisitions — year & price\"",
    "Demo 3 (Practical): \"Dinner for 8 on $100 — 2 vegetarian, 1 gluten-free\"",
    "Watch for: tone differences, confidence levels, formatting, accuracy",
    "Key question: Which response would you actually USE?",
],
notes="Same prompt typed into all three. ChatGPT tends to be dramatic/punchy, Claude more literary/careful, Gemini concise. For factual questions, notice who hedges vs states as fact. For practical tasks, notice who gives something actionable."
)

# Slide: Brain vs App
add_content_slide("The Brain vs The App", [
    "The App = the website you talk to (ChatGPT, Claude.ai, Gemini)",
    "The Brain = the AI model inside (GPT-5.2, Claude Sonnet 4.6, Gemini 3 Pro)",
    "Netflix analogy: Netflix is the app, shows are the brains",
    "One app can run different brains (ChatGPT offers GPT-5.2, o3, GPT-4o…)",
    "The SAME brain can run in different apps (Claude via claude.ai OR OpenClaw)",
    "OpenClaw = Netflix + Disney+ + HBO in one app",
],
notes="This is the most important concept. When you open ChatGPT you're using ONE app with ONE company's brains. That's like only watching Netflix. By end of today you'll pick the best brain for each job regardless of which app it lives in."
)

# Slide: Brain vs App table
add_table_slide("The Brain vs The App — Cheat Sheet",
    ["App (The Netflix)", "Brains Available (The Shows)"],
    [
        ["ChatGPT", "GPT-5.2, GPT-5.1, GPT-4o, o3, o4-mini"],
        ["Claude.ai", "Claude Opus 4.6, Sonnet 4.6, Haiku 4.5"],
        ["Gemini", "Gemini 3 Pro, Gemini 2.0 Flash"],
        ["OpenClaw", "ALL of the above — one interface, every brain"],
    ],
notes="See that last row? That's the goal today. One tool that lets you pick any brain."
)

# ══ PART 2: HOW IT WORKS ═════════════════════════════════════════════════════

add_section_slide("Part 2: How It Works", "Models, pricing, and API keys — the stuff nobody explains clearly")

# Slide: Model tiers
add_table_slide("Models — Fast/Cheap vs Smart/Expensive",
    ["Tier", "Examples (Feb 2026)", "Speed", "Cost/msg", "Best For"],
    [
        ["🐇 Fast & Cheap", "Haiku 4.5, Gemini Flash, GPT-4o-mini", "~1-2s", "~$0.001", "Quick questions, brainstorming"],
        ["🦊 Balanced", "Sonnet 4.6, GPT-5.2, Gemini 3 Pro", "~3-5s", "~$0.01", "Most real work, writing, analysis"],
        ["🧠 Max Power", "Opus 4.6, o3", "~10-15s", "~$0.05", "Complex reasoning, deep analysis"],
    ],
    col_widths=[1.5, 3.2, 1.2, 1.3, 4.3],
notes="Cars analogy: Haiku = Honda Civic (errands), Sonnet = BMW (road trip), Opus = F1 car (the race). The expensive model isn't always better — 'What's the capital of France?' gets the same answer from all three at 50x cost difference."
)

# Slide: Feb 2026 model landscape
add_table_slide("The AI Model Landscape — February 2026",
    ["Company", "Flagship", "Mid-Tier", "Fast/Cheap", "Reasoning"],
    [
        ["OpenAI", "GPT-5.2 (Dec '25)\n400K ctx, $20/$60", "GPT-5.1, GPT-4o", "GPT-4o-mini", "o3, o4-mini"],
        ["Anthropic", "Opus 4.6 (Feb '26)", "Sonnet 4.6 (Feb '26)", "Haiku 4.5 (Oct '25)", "—"],
        ["Google", "Gemini 3 Pro (Nov '25)\n1M ctx", "—", "Gemini 2.0 Flash\n650ms avg!", "—"],
        ["Others", "Grok 4 (xAI)", "Mistral Large", "DeepSeek V4\n(94% cheaper)", "Llama 4 (open)"],
    ],
    col_widths=[1.5, 3.0, 2.5, 2.5, 2.0],
notes="This landscape changes fast. Key takeaway: competition is fierce, prices dropping, capabilities rising. As of Feb 2026 these are the latest flagships."
)

# Slide: Free vs Paid
add_table_slide("Free Tiers vs Paid — The Pricing Truth",
    ["Service", "Free Tier", "Paid Plan", "API (OpenClaw)"],
    [
        ["ChatGPT", "GPT-4o-mini unlimited\nGPT-4o limited (~15 msg)", "$20/mo → more GPT-5.2, o3", "Pay per message\n~$3-5/mo light use"],
        ["Claude.ai", "Sonnet 4.6 limited\n(~20-30 msg then locked)", "$20/mo → 5x more,\nOpus 4.6 access", "Pay per message\n~$3-5/mo light use"],
        ["Gemini", "Gemini 3 Pro\n(generous free tier)", "$20/mo → longer context,\nGoogle integration", "Generous free tier\nor pennies per msg"],
    ],
    col_widths=[1.5, 3.0, 3.0, 4.0],
notes="Free = casino (happy for you to walk in but there's a catch). $20/mo = one company. API = pay-as-you-go for ALL companies. Most people spend $5-15/mo via API."
)

# Slide: API Key
add_content_slide("What's an API Key?", [
    "Analogy: a prepaid gas card for AI",
    "Go to OpenAI/Anthropic/Google → create account → add $10 credit",
    "They give you a long code: sk-proj-abc123def456...",
    "Paste it into OpenClaw → every message deducts pennies",
    "Treat it like a credit card number — never share it publicly",
    "You set it up once and forget about it (like Wi-Fi password)",
],
notes="Prepaid gas card analogy. Load money, get a code, use gas, balance goes down. sk-proj-abc123... is just proof you've paid. If someone gets it they can use your balance. We'll set one up together in Part 6."
)

# Slide: OAuth vs API Key
add_table_slide("OAuth vs API Key — Two Ways to Log In",
    ["", "OAuth (\"Sign in with...\")", "API Key"],
    [
        ["Analogy", "Hotel key card\n(front desk verifies you)", "Prepaid gas card\n(loaded money, use anytime)"],
        ["How it feels", "Click \"Sign in with Google\"", "Copy-paste a long code"],
        ["Who tracks cost", "App (included in subscription)", "You (pay-per-use from balance)"],
        ["Example", "Using ChatGPT's website", "Using OpenClaw with your key"],
        ["Flexibility", "Stuck with that app's options", "Use any tool that accepts the key"],
    ],
    col_widths=[2.0, 4.5, 5.0],
notes="OAuth = hotel key card, works at that hotel only. API key = prepaid gas card, works at any pump. You already use OAuth daily. We'll add API keys in Part 6."
)

# ══ PART 3: WHAT AI CAN DO ═══════════════════════════════════════════════════

add_section_slide("Part 3: What AI Can Do", "Image generation, voice conversations, and vision — live demos")

# Slide: Image Gen
add_content_slide("Image Generation — AI Creates Pictures from Words", [
    "Describe what you want in plain English → AI draws it in ~10 seconds",
    "Demo: \"A watercolor golden retriever in a business suit at a desk\"",
    "ChatGPT has DALL-E built in; Gemini uses Imagen; Claude doesn't generate images",
    "Great for: social media graphics, presentations, mockups, brainstorming",
    "Not perfect: check hands/fingers, watch for weird artifacts",
    "Practical, not portfolio-quality — \"good enough\" for most needs",
],
notes="Live demo: type the golden retriever prompt. Show iteration ('now make him giving a TED talk to cats'). Point out it understood style, spatial arrangement, humor."
)

# Slide: Voice
add_content_slide("Voice — Talk to AI Like a Phone Call", [
    "ChatGPT app → tap headphone icon → full voice conversation",
    "Sounds natural, not robotic. Remembers context within conversation.",
    "Demo: \"Explain AI to my 75-year-old mother-in-law who fears robots\"",
    "Great for: hands-free (driving/cooking), accessibility, brainstorming on walks",
    "This isn't Siri giving Wikipedia snippets — it's real back-and-forth",
],
notes="Do this on phone with speakers. Let the audience hear the response. Follow up with 'now make it funnier'. Show that it maintains context."
)

# Slide: Vision
add_content_slide("Vision — AI Can See and Understand Images", [
    "Upload any photo → AI identifies objects, reads text, understands context",
    "Demo 1: Photo of the room → \"What do you see? Be specific.\"",
    "Demo 2: Restaurant receipt → \"Split 4 ways, 2 didn't drink, 18% tip\"",
    "Demo 3: Error message screenshot → \"What does this mean? Fix?\"",
    "Claude excels at vision; ChatGPT & Gemini also strong",
    "Plant ID, handwriting OCR, math homework, foreign menus — all work",
],
notes="Vision turns AI from 'a thing I type to' into 'a thing that can look at my world.' These three capabilities separate 'I've tried ChatGPT' from 'I actually use AI.'"
)

# ══ PART 4: WHAT AI GETS WRONG ═══════════════════════════════════════════════

add_section_slide("Part 4: What AI Gets Wrong", "Hallucinations and privacy — the important warnings")

# Slide: Hallucinations
add_content_slide("Hallucinations — When AI Is Confidently Wrong", [
    "AI makes things up — regularly, confidently, with perfect grammar",
    "Demo: Ask about fake case \"Robertson v. McKenzie (1987)\" → AI invents details",
    "Real case: NY lawyer sanctioned for citing AI-generated fake cases",
    "Not lying (no intentions) — more like a confident student who didn't study",
    "HIGH risk: specific dates, names, statistics, niche domains",
    "LOW risk: general knowledge, creative tasks, reformatting YOUR content",
],
notes="This case doesn't exist — I made it up. But the AI will write a detailed, convincing response. A journalist would publish it. A lawyer would cite it. 100% fiction."
)

# Slide: How to verify
add_content_slide("The 30-Second Verification Rule", [
    "🔍  Google the specific claim (10 seconds)",
    "🔄  Ask a second AI the same question (if they disagree, dig deeper)",
    "❓  Ask the AI: \"Are you sure? Can you provide a checkable source?\"",
    "📋  For anything important: verify with a real, primary source",
    "Rule: AI is your first draft, never your final answer",
    "AI is amazing at thinking, mediocre at remembering",
],
notes="Use AI for brainstorming, drafting, getting started. If the fact matters — report, email to client, school paper — spend 30 seconds verifying."
)

# Slide: Privacy
add_table_slide("Privacy & Data — Where Does Your Stuff Go?",
    ["Service", "Training on Your Data?", "Notes"],
    [
        ["ChatGPT (free)", "Yes by default", "Opt out available in Data Controls"],
        ["ChatGPT Plus", "Yes by default", "Toggle off in settings"],
        ["Claude.ai", "No (by default)", "Anthropic doesn't train on conversations"],
        ["Gemini", "Yes", "Integrated with your Google account"],
        ["API (OpenClaw)", "No", "Paying customers — not free training data"],
    ],
notes="Free ChatGPT = your conversations can train future models. API route is safest. Treat AI like a smart coworker you don't fully trust yet."
)

# Slide: Privacy rules
add_content_slide("Privacy Rules of Thumb", [
    "✅ SAFE: Public info, your creative writing, general questions",
    "⚠️ THINK TWICE: Client details, internal strategies, financial info, NDAs",
    "❌ NEVER: Passwords, credit card/SIN numbers, medical records with names",
    "Use the API for sensitive work (companies don't train on API traffic)",
    "Check settings: ChatGPT → Data Controls → toggle off training",
    "Bottom line: treat AI like a smart coworker you don't fully trust yet",
],
notes="Never paste passwords or personal health info. Be cautious with confidential docs on free tiers. API is the safest option."
)

# Break slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, DARK_BG)
tb = _tb(slide, 1.0, 2.5, 11, 1.5)
_set(tb.text_frame, "☕  10-Minute Break", size=44, bold=True, color=WHITE)
tb2 = _tb(slide, 1.0, 4.0, 11, 1.0)
_set(tb2.text_frame, "Next up: Hands-on with all three AIs", size=24, color=ACCENT)

# ══ PART 5: HANDS-ON ═════════════════════════════════════════════════════════

add_section_slide("Part 5: Hands-On", "Time to get your hands dirty — 40 minutes of practice")

# Slide: Exercise 1
add_content_slide("Exercise 1: Sign Up & First Message (10 min)", [
    "Open claude.ai → Sign up (Google account or email)",
    "Open gemini.google.com → Sign in with Google",
    "First prompt: \"I'm new here. Tell me one thing most people don't know you can do.\"",
    "Second prompt: \"I have a job interview for marketing manager at a tech company…\"",
    "Compare: Different tone? Detail level? Did one ask follow-up questions?",
],
notes="Walk around the room. Most common issue: people not finding the sign-up button. Claude needs email verification. Gemini usually auto-signs in with Google."
)

# Slide: Exercise 2
add_content_slide("Exercise 2: The AI Taste Test (15 min)", [
    "Pick ONE task — paste the SAME prompt into ChatGPT, Claude, AND Gemini:",
    "Option A: \"Write a professional email to Sarah — project 2 weeks late…\"",
    "Option B: \"Lease vs buy a car — 20K km/yr, 5+ years, $500/mo budget…\"",
    "Option C: \"Explain mortgage rates to a first-time buyer — fixed vs variable…\"",
    "Compare: Which was most helpful? Longest ≠ best? Who asked clarifying Qs?",
    "Develop YOUR taste — you'll start to notice preferences",
],
notes="Don't rush. Actually read all three. The point isn't to find a winner — it's to develop your own intuition."
)

# Slide: Exercise 3
add_content_slide("Exercise 3: Upload a Document + Image (15 min)", [
    "Part A: Upload a PDF/Word doc to Claude → \"Summarize in 5 bullet points\"",
    "Part B: Take a photo with your phone of ANYTHING nearby",
    "Upload to Claude or ChatGPT → \"What's in this image? Describe everything.\"",
    "Try: receipt splitting, error message diagnosis, plant identification",
    "This is the transition from \"AI is fun\" to \"AI saves me time\"",
],
notes="This is where it clicks. When AI reads your 20-page report or your terrible handwriting, it stops being a toy and starts being a tool."
)

# ══ PART 6: OPENCLAW & SETUP ═════════════════════════════════════════════════

add_section_slide("Part 6: Getting Ready for OpenClaw", "Markdown, Linux basics, and setting up your personal AI assistant")

# Slide: What is Markdown?
add_content_slide("What is Markdown? (And Why AI Tools Love It)", [
    "Plain text with simple formatting: # heading, **bold**, - bullet",
    "No special software needed — any text editor works",
    "AI tools speak Markdown natively (prompts, responses, configs)",
    "GitHub, Discord, Slack, Notion all use Markdown",
    "Example: **bold** = bold,  _italic_ = italic,  # = heading",
    "OpenClaw config files, memory, and agents all use Markdown",
],
notes="Markdown is the universal language of AI tools. If you learn one syntax this year, make it this. It takes 5 minutes to learn the basics."
)

# Slide: Markdown examples
add_table_slide("Markdown Quick Reference",
    ["You Type", "You Get", "Use Case"],
    [
        ["# Heading", "Large heading", "Document sections"],
        ["**bold text**", "bold text", "Emphasis"],
        ["- item", "• bullet point", "Lists"],
        ["1. item", "1. numbered item", "Ordered lists"],
        ["`code`", "inline code", "Commands, file names"],
        ["[text](url)", "clickable link", "References"],
    ],
notes="These six things cover 90% of what you'll need. Practice in Discord — it renders Markdown natively."
)

# Slide: Linux Survival Kit
add_content_slide("Linux Survival Kit — Terminal Basics", [
    "OpenClaw runs in a terminal (WSL on Windows, Terminal on Mac)",
    "cd folder — change directory (like double-clicking a folder)",
    "ls — list files (like opening a folder to see what's inside)",
    "cat file.txt — display file contents (like opening a document)",
    "nano file.txt — edit a file (like Notepad, but in terminal)",
    "Ctrl+C — stop/cancel anything that's running",
],
notes="You don't need to become a Linux expert. These 5 commands cover 90% of what you'll need for OpenClaw. cd = go somewhere, ls = look around, cat = read, nano = edit, Ctrl+C = panic button."
)

# Slide: Linux cheat sheet
add_table_slide("Linux Terminal Cheat Sheet",
    ["Command", "What It Does", "Analogy"],
    [
        ["cd documents", "Go into the documents folder", "Double-click a folder"],
        ["cd ..", "Go back one folder", "Click the back button"],
        ["ls", "List everything in current folder", "Open folder to see files"],
        ["cat notes.txt", "Show contents of notes.txt", "Open a document to read it"],
        ["nano notes.txt", "Edit notes.txt", "Open in Notepad"],
        ["pwd", "Print where you are", "Check folder path in address bar"],
        ["Ctrl+C", "Stop whatever's running", "Force-quit"],
    ],
notes="Keep this cheat sheet handy. You'll memorize these within a week of using OpenClaw."
)

# Slide: Meet OpenClaw
add_content_slide("Meet OpenClaw — One App, Every AI Brain", [
    "Personal AI assistant: ANY model through one interface",
    "Runs on phone (Discord), computer (CLI), or web",
    "Instead of $20/mo to each company → pennies per message via API",
    "Live demo: same question → switch model → different brain, same app",
    "That conversation cost ~2 cents. Not $20/month.",
],
notes="Remember the Netflix analogy? OpenClaw = Netflix + Disney+ + HBO in one app. Show Discord on phone, type a message, switch models with /model, show cost."
)

# Slide: Cost comparison
add_table_slide("Real Cost Comparison",
    ["Usage", "ChatGPT Plus", "Claude Pro", "OpenClaw (API)"],
    [
        ["Light (few times/week)", "$20/mo (overpaying)", "$20/mo (overpaying)", "~$1-2/mo"],
        ["Moderate (daily use)", "$20/mo (one company)", "$20/mo (one company)", "~$5-15/mo (ALL companies)"],
        ["Heavy (power user)", "$20/mo + $20/mo = $40", "$20/mo + $20/mo = $40", "~$15-25/mo (everything)"],
    ],
    col_widths=[2.5, 3.0, 3.0, 3.0],
notes="Smart model switching: quick question → Haiku ($0.0005), real work → Sonnet ($0.01), deep analysis → Opus ($0.05). Match the brain to the task."
)

# Slide: Smart model switching demo
add_content_slide("Smart Model Switching — Demo", [
    "🐇 Quick Q → claude-haiku-4.5: \"Substitute for buttermilk?\" → $0.0005",
    "🦊 Real work → claude-sonnet-4.6: \"Make this email professional\" → $0.01",
    "🧠 Deep analysis → claude-opus-4.6: \"Toronto vs Calgary financial analysis\" → $0.05",
    "Like cars: don't take a taxi to the mailbox, don't bike to the airport",
    "OpenClaw makes switching instant: /model [name]",
],
notes="Live demo each tier. Show the cost after each response. The pattern clicks when people see the 100x price difference for a buttermilk question."
)

# Slide: Get API keys
add_content_slide("Step 1: Get Your API Keys (10 min)", [
    "OpenAI: platform.openai.com → API Keys → Create → add $10 credit",
    "Anthropic: console.anthropic.com → API Keys → Create → add $10 credit",
    "Google: aistudio.google.com → Get API Key → Create (generous free tier)",
    "Name each key \"openclaw\" — copy it immediately (you can't see it again!)",
    "Treat keys like credit card numbers — never share publicly",
    "$10 per provider will last weeks of normal use",
],
notes="I know this feels like a lot of copying and pasting. You do it once and you're done forever. Like setting up Wi-Fi — annoying for 10 minutes, never think about it again."
)

# Slide: Install — Windows
add_content_slide("Step 2a: Windows Setup (WSL + OpenClaw)", [
    "1. PowerShell (Admin): wsl --install → restart computer",
    "2. Open Ubuntu from Start menu → create username/password",
    "3. sudo apt update && sudo apt upgrade -y",
    "4. curl -fsSL https://deb.nodesource.com/setup_22.x | sudo -E bash -",
    "5. sudo apt install -y nodejs",
    "6. npm install -g @anthropic/openclaw && openclaw setup",
],
notes="WSL = Windows Subsystem for Linux. A little Linux computer inside your Windows computer. Most annoying part. Once done, everything else is easy. Password won't show dots — that's normal."
)

# Slide: Install — Mac
add_content_slide("Step 2b: macOS Setup", [
    "1. Open Terminal (Cmd+Space → \"Terminal\" → Enter)",
    "2. Install Homebrew: /bin/bash -c \"$(curl -fsSL ...)\"",
    "3. brew install node",
    "4. npm install -g @anthropic/openclaw",
    "5. openclaw setup → paste your API keys when prompted",
    "6. Test: openclaw chat \"Hello! Tell me a fun fact.\"",
],
notes="macOS is simpler — no WSL needed. Homebrew is a package manager, like an App Store for terminal tools."
)

# Slide: Verify
add_content_slide("Step 3: Verify It Works", [
    "In your terminal: openclaw chat \"Hello! Tell me a fun fact.\"",
    "If you see a response → you're done! Every major AI brain is available.",
    "Common issues: \"command not found\" → run node --version to check",
    "\"invalid API key\" → go back to platform, create a new one",
    "\"insufficient funds\" → add credit to your API account",
    "I'll stay after the session to help anyone who's stuck",
],
notes="Hands up if you got a response! Celebrate. Most common issue is confusing ChatGPT login with API platform login — they're separate."
)

# ── CLOSING ──────────────────────────────────────────────────────────────────

add_content_slide("What to Remember from Today", [
    "There are 3+ major AIs — ChatGPT, Claude, Gemini — each with strengths",
    "The Brain ≠ The App — pick the best model for each task",
    "AI hallucinates — always verify facts (30-second rule)",
    "API keys give you freedom & save money vs subscriptions",
    "OpenClaw = one interface for every AI model, pay-per-use",
    "Start with Claude Sonnet 4.6 for most things — best all-rounder in Feb 2026",
],
notes="You walked in knowing ChatGPT. You're walking out knowing how to use ANY AI, switch between them, and pay a fraction of subscription costs."
)

# Slide: Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, DARK_BG)
tb = _tb(slide, 1.0, 2.0, 11, 1.5)
_set(tb.text_frame, "Questions & Answers", size=48, bold=True, color=WHITE)
tb2 = _tb(slide, 1.0, 3.8, 11, 2.5)
tf = tb2.text_frame
_set(tf, "Common Qs:", size=22, bold=True, color=ACCENT)
for q in [
    "\"Will AI take my job?\" — No, but someone who uses AI well might outperform you.",
    "\"Which AI daily?\" — Claude Sonnet 4.6 for most things. GPT-5.2 for images. Gemini for Google integration.",
    "\"Is my data safe?\" — API is safest. Never put passwords in any AI.",
    "\"Open source?\" — Llama 4 is great, DeepSeek V4 is 94% cheaper. Advanced session topic.",
]:
    _add_para(tf, q, size=18, color=WHITE, space_before=Pt(14))

slide.notes_slide.notes_text_frame.text = "Nothing is too basic — if you're wondering it, three other people are too."

# ── SAVE ─────────────────────────────────────────────────────────────────────
out = "/home/lj_wsl/cut-the-crap/general/cut-the-crap-general.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides)} slides")
