#!/usr/bin/env python3
"""Generate Cut the Crap PowerPoint deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# Colors
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)  # Dark blue
ACCENT = RGBColor(0x00, 0x96, 0xD6)    # Bright blue accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xFE)
CHARCOAL = RGBColor(0x23, 0x23, 0x2E)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_dark_slide(title_text, subtitle_text="", notes=""):
    """Dark background title/section slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(22)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3), Inches(2.333), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_content_slide(title_text, bullets=None, notes="", two_col=False):
    """White background content slide with bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BG
    bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if bullets:
        if two_col and len(bullets) > 5:
            mid = len(bullets) // 2
            col1 = bullets[:mid]
            col2 = bullets[mid:]
            _add_bullet_box(slide, col1, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.5))
            _add_bullet_box(slide, col2, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.5))
        else:
            _add_bullet_box(slide, bullets, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def _add_bullet_box(slide, bullets, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Support indented bullets (start with "  ")
        if isinstance(bullet, tuple):
            p.text = bullet[1]
            p.level = bullet[0]
        else:
            p.text = bullet
            p.level = 0
        
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
        p.space_before = Pt(4)


def add_table_slide(title_text, headers, rows, notes=""):
    """Slide with a styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BG
    bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table_width = Inches(11.5)
    row_height = Inches(0.55)
    table_height = row_height * num_rows
    
    tbl_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.9), Inches(1.5), table_width, table_height)
    tbl = tbl_shape.table
    
    # Set column widths evenly
    col_w = int(table_width / num_cols)
    for i in range(num_cols):
        tbl.columns[i].width = col_w
    
    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
    
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = DARK_TEXT
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_diagram_slide(title_text, lines, notes=""):
    """Slide with monospaced text block (for diagrams)."""
    slide = add_content_slide(title_text, notes=notes)
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.name = "Consolas"
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(2)
    return slide


# ============================================================
# BUILD THE DECK
# ============================================================

# --- TITLE SLIDE ---
add_dark_slide(
    "Cut the Crap",
    "Everything You Need to Know About AI"
)

# ============================================================
# PART 1: THE LANDSCAPE
# ============================================================
add_dark_slide("Part 1", "The Landscape")

# Slide: There's More Than ChatGPT
add_content_slide(
    "There's More Than ChatGPT",
    [
        "ChatGPT (OpenAI) — The one everyone knows. Great all-rounder, best ecosystem",
        "Claude (Anthropic) — More careful, better at long documents, less likely to make stuff up",
        "Gemini (Google) — Integrated with Google services, good at research, handles images well",
        "",
        "None of these is \"the best\" — they're like different restaurants",
        "By the end of today you'll know when to use which one"
    ],
    notes="Raise your hand if you've used ChatGPT. Now keep your hand up if you've also tried Claude or Gemini. That's totally normal. ChatGPT had a massive head start — it launched November 2022 and became the fastest-growing app in history. But there are now three serious players, and they're all good at different things."
)

# Slide: Live Comparison
add_content_slide(
    "Live Comparison — Same Prompt, 3 AIs",
    [
        "Demo 1 — Creative Writing: See personality differences",
        "  → \"Write a 4-sentence horror story about a smart home that becomes too smart\"",
        "",
        "Demo 2 — Factual Research: See accuracy differences",
        "  → \"What were the three biggest acquisitions in Canadian tech history?\"",
        "",
        "Demo 3 — Practical Task: See usefulness differences",
        "  → Dinner party for 8 on $100 budget (2 vegetarian, 1 gluten-free)"
    ],
    notes="Same prompt, three different personalities. None is wrong — they just have different vibes. Notice which AI is confident and which is cautious. Cautious is sometimes better."
)

# Slide: The Brain vs The App
add_content_slide(
    "The Brain vs The App",
    [
        "THE APP = the website you talk to (ChatGPT, Claude, Gemini)",
        "THE BRAIN = the AI model inside (GPT-4o, Claude Sonnet, Gemini Pro)",
        "",
        "Think of Netflix vs the shows on Netflix:",
        "  • Netflix = delivery system (the app)",
        "  • Shows = the actual content (the brain)",
        "",
        "One app can run DIFFERENT brains",
        "The SAME brain can run in DIFFERENT apps"
    ],
    notes="This is the single most important concept today. ChatGPT is like Netflix — it's the app. GPT-4o is like a show on Netflix — it's the actual AI brain doing the thinking. When you open ChatGPT, you're using ONE app with ONE company's brains. By the end of today, you'll be able to pick the best brain for each job."
)

# Slide: Brain vs App Table
add_table_slide(
    "Apps and Their Brains",
    ["App (The Netflix)", "Brains (The Shows)"],
    [
        ["ChatGPT", "GPT-4o, GPT-4o-mini, o1, o3"],
        ["Claude.ai", "Claude Sonnet, Claude Haiku, Claude Opus"],
        ["Gemini", "Gemini Pro, Gemini Flash"],
        ["OpenClaw", "All of the above ← this is where we're headed"],
    ],
    notes="See that last row? That's the goal today. One tool that lets you pick any brain. Like having Netflix, Disney+, and HBO all in one app."
)

# ============================================================
# PART 2: HOW IT WORKS
# ============================================================
add_dark_slide("Part 2", "How It Works")

# Slide: Models — Fast/Cheap vs Smart/Expensive
add_table_slide(
    "Models — Fast/Cheap vs Smart/Expensive",
    ["Tier", "Examples", "Speed", "Cost", "Good For"],
    [
        ["🐇 Fast & Cheap", "GPT-4o-mini, Haiku, Flash", "~2 sec", "~$0.001/msg", "Quick questions, brainstorming"],
        ["🦊 Balanced", "GPT-4o, Sonnet, Pro", "~5 sec", "~$0.01/msg", "Most real work, writing, analysis"],
        ["🧠 Max Power", "o3, Claude Opus", "~15 sec", "~$0.05/msg", "Complex reasoning, deep analysis"],
    ],
    notes="Think of it like cars: Haiku/Flash/Mini = Honda Civic. Sonnet/GPT-4o = BMW. Opus/o3 = Formula 1 car. The expensive model isn't always better. If you ask 'What's the capital of France?' the Civic and the F1 car give the same answer. The skill is matching the task to the model."
)

# Slide: Free Tiers vs Paid
add_table_slide(
    "Free Tiers — What You Actually Get",
    ["Service", "Free Tier", "Limit"],
    [
        ["ChatGPT", "GPT-4o-mini unlimited, GPT-4o limited", "~10-15 GPT-4o msgs, then throttled"],
        ["Claude.ai", "Claude Sonnet, limited", "~20-30 msgs, then locked for hours"],
        ["Gemini", "Gemini Pro, generous", "Rarely hit limits for normal use"],
    ],
    notes="Free in AI is like free at a casino — they're happy for you to walk in, but there's a catch."
)

add_table_slide(
    "Paid Plans vs API (Pay-Per-Use)",
    ["Option", "Price", "What You Get"],
    [
        ["ChatGPT Plus", "$20/mo", "More GPT-4o, o1 access, image gen"],
        ["Claude Pro", "$20/mo", "5x more Sonnet, Opus access"],
        ["Gemini Advanced", "$20/mo", "Gemini Ultra, Google integration"],
        ["API via OpenClaw (light)", "$1–5/mo", "All brains, pay only what you use"],
        ["API via OpenClaw (moderate)", "$5–15/mo", "All brains, daily use"],
        ["API via OpenClaw (heavy)", "$15–40/mo", "All brains, power user"],
    ],
    notes="$20/month plans lock you to ONE company. The API route via OpenClaw gives access to ALL brains and most people spend $5-15/month."
)

# Slide: What's an API Key?
add_content_slide(
    "What's an API Key?",
    [
        "Think of a prepaid gas card:",
        "  1. Go to OpenAI's website, create an account",
        "  2. Add $10 to your balance (like loading a gift card)",
        "  3. They give you a long code — that's your API key",
        "  4. Paste that code into OpenClaw",
        "  5. Every message deducts a tiny amount from your balance",
        "",
        "⚠️  Treat it like a credit card number — don't share it!",
        "",
        "sk-proj-abc123def456ghi789..."
    ],
    notes="An API key is your prepaid gas card for AI. You load money, use AI, done. The key is just proof that you've paid. If someone gets your API key, they can use your balance. We'll set one up together in Part 6."
)

# Slide: OAuth vs API Key
add_table_slide(
    "OAuth vs API Key — Two Ways to Log In",
    ["", "OAuth (\"Sign in with...\")", "API Key"],
    [
        ["Analogy", "Hotel key card", "Prepaid gas card"],
        ["How it feels", "Click 'Sign in with Google' → done", "Copy-paste a long code into settings"],
        ["Who tracks cost", "The app (subscription)", "You (pay-per-use)"],
        ["Example", "Using ChatGPT's website", "Using OpenClaw with your own key"],
        ["Flexibility", "Stuck with that app's options", "Use any tool that accepts the key"],
    ],
    notes="OAuth = simple login on one website. API key = portable pass you can use anywhere. Website login vs portable code. That's it."
)

# ============================================================
# PART 3: WHAT AI CAN DO
# ============================================================
add_dark_slide("Part 3", "What AI Can Do")

# Image Generation
add_content_slide(
    "AI Can Create Images from Words",
    [
        "Describe what you want → AI draws it in ~10 seconds",
        "",
        "Demo: \"Create a watercolor painting of a golden retriever wearing",
        "  a tiny business suit, sitting at a desk, looking very serious\"",
        "",
        "It understands: style, objects, spatial arrangement, concepts",
        "It can iterate: \"Now make the same dog giving a TED talk to cats\"",
        "",
        "Best for: social media graphics, presentations, mockups",
        "Not perfect: check hands, fingers, text in images",
        "",
        "ChatGPT → DALL-E (built in)  |  Gemini → Imagen  |  Claude → not yet"
    ],
    notes="Image gen is great for 'good enough' visuals. It's not replacing graphic designers for polished work, but for a quick social post, a presentation image, or a concept mockup? Incredible."
)

# Voice
add_content_slide(
    "You Can TALK to AI — Like a Phone Call",
    [
        "ChatGPT voice mode: tap the headphone icon",
        "",
        "Demo: \"I'm trying to explain to my mother-in-law what AI is.",
        "  She's 75 and thinks robots are taking over.",
        "  Give me a simple, reassuring way to explain it.\"",
        "",
        "It sounds natural, not robotic",
        "It remembers context within the conversation",
        "",
        "Great for: hands-free (driving/cooking), accessibility,",
        "  language practice, brainstorming while walking",
        "",
        "This isn't Siri giving you a Wikipedia snippet —",
        "  this is a real back-and-forth conversation"
    ],
    notes="This isn't Siri giving you a Wikipedia snippet. This is a real back-and-forth conversation. My favorite use: when I'm driving and need to think through a problem, I just talk to it."
)

# Vision
add_content_slide(
    "AI Can See — Upload a Photo",
    [
        "Upload a photo → AI understands what's in it",
        "",
        "Demo 1: Upload photo of something in the room",
        "  → \"What do you see in this image? Be specific.\"",
        "",
        "Demo 2: Upload a restaurant receipt",
        "  → \"Split this bill 4 ways. Two people didn't drink. Include 18% tip.\"",
        "",
        "Demo 3: Upload an error message screenshot",
        "  → \"What does this error mean? Explain like I'm not technical.\"",
        "",
        "Also works for: plant identification, handwriting reading,",
        "  diagram analysis, foreign language menus, homework help"
    ],
    notes="Vision turns AI from 'a thing I type to' into 'a thing that can look at my world.' These three capabilities — image gen, voice, vision — are what separate 'I've tried ChatGPT' from 'I actually use AI.'"
)

# ============================================================
# PART 4: WHAT AI GETS WRONG
# ============================================================
add_dark_slide("Part 4", "What AI Gets Wrong")

# Hallucinations
add_content_slide(
    "Hallucinations — When AI Is Confidently Wrong",
    [
        "AI makes things up. Regularly. With complete confidence.",
        "",
        "It's not lying — it's like a confident student who didn't study:",
        "  instead of saying \"I don't know,\" it generates something that SOUNDS right",
        "",
        "Real case: A New York lawyer submitted a brief with",
        "  AI-generated case citations. The cases were fake.",
        "  He was sanctioned by the judge. Front-page news.",
        "",
        "Demo: Ask about \"Robertson v. McKenzie (1987)\"",
        "  → AI writes a detailed, convincing response",
        "  → This case does not exist. It was made up."
    ],
    notes="AI makes things up. Not occasionally. Regularly. And it does it with complete confidence, perfect grammar, and a straight face. This is called a 'hallucination.'"
)

# Verification
add_content_slide(
    "The 30-Second Verification Rule",
    [
        "Before trusting AI on facts, do ONE of these:",
        "  🔍 Google the specific claim (takes 10 seconds)",
        "  🔄 Ask a second AI the same question",
        "  ❓ Ask the AI: \"Are you sure? Can you provide a source?\"",
        "  📋 For anything important: verify with a real source",
        "",
        "HIGH hallucination risk: dates, numbers, names, citations, niche topics",
        "LOW hallucination risk: general knowledge, creative tasks, formatting",
        "",
        "Rule: AI is your first draft, never your final answer",
        "",
        "AI is amazing at thinking, mediocre at remembering"
    ],
    notes="Use it for brainstorming, drafting, getting started. But if the fact matters — if it's going in a report, an email to a client, a school paper — spend 30 seconds verifying."
)

# Privacy
add_table_slide(
    "Privacy & Data — Where Does Your Stuff Go?",
    ["What you type/upload", "What happens to it"],
    [
        ["ChatGPT (free)", "OpenAI CAN use it to train future models (opt out available)"],
        ["ChatGPT Plus", "Same default, but you can toggle off in settings"],
        ["Claude.ai", "Anthropic does NOT use conversations for training (by default)"],
        ["Gemini", "Google CAN use it, integrated with your Google account"],
        ["API (via OpenClaw)", "Generally NOT used for training — you're a paying customer"],
    ],
    notes="By default, free ChatGPT conversations can be used to train future AI models. Use the API for sensitive work. Check the 'Data Controls' settings."
)

add_content_slide(
    "What's Safe to Put in AI?",
    [
        "✅ SAFE: Public info, your own writing, general questions, anonymized scenarios",
        "",
        "⚠️ THINK TWICE: Client names/details, internal business strategies,",
        "  personal financial info, anything under NDA",
        "",
        "❌ NEVER: Passwords, credit card / SIN numbers, medical records",
        "  with names, other people's private info without consent",
        "",
        "Bottom line: treat AI like a smart coworker you don't fully trust yet",
        "  — brainstorm with them, but don't hand them your bank password"
    ],
    notes="Treat AI like a smart coworker you don't fully trust yet. You'd brainstorm with them, ask for help, run ideas by them — but you wouldn't hand them your bank password or your client's confidential files."
)

# ============================================================
# PART 5: HANDS-ON
# ============================================================
add_dark_slide("Part 5", "Hands-On Exercises")

add_content_slide(
    "Exercise 1: Sign Up & First Message",
    [
        "Go to claude.ai — Sign up (Google account or email)",
        "Go to gemini.google.com — Sign in with Google",
        "",
        "Try this prompt in both:",
        "  \"I'm new here. Tell me one thing you can do",
        "   that most people don't know about.\"",
        "",
        "Then try:",
        "  \"I have a job interview next week for a marketing manager",
        "   position at a tech company. Give me 5 questions they're likely",
        "   to ask and help me draft strong answers.\"",
        "",
        "Compare: Different tone? Detail level? Which would you use?"
    ],
    notes="Walk around the room, help people sign up. Most common issues: email verification, choosing the right Google account."
)

add_content_slide(
    "Exercise 2: The AI Taste Test",
    [
        "Pick ONE task. Paste into ChatGPT, Claude, AND Gemini:",
        "",
        "Option A — Email Writer:",
        "  Write a professional email to Sarah explaining project is 2 weeks late",
        "",
        "Option B — Decision Helper:",
        "  Leasing vs buying a car (20k km/year, 5+ years, $500/mo budget)",
        "",
        "Option C — Learning:",
        "  Explain mortgage interest rates to a first-time home buyer",
        "",
        "Compare: Most helpful? Best length? Did any ask clarifying questions?"
    ],
    notes="Don't rush this. Actually read all three. The point isn't to find a winner — it's to develop your own taste."
)

add_content_slide(
    "Exercise 3: Upload a Document & Image",
    [
        "Part A — Upload a Document to Claude:",
        "  Click 📎 → select a PDF, Word doc, or paste article text",
        "  → \"Summarize in 5 bullet points. What are the key takeaways?\"",
        "",
        "Part B — Upload an Image:",
        "  Take a photo of ANYTHING: whiteboard, receipt, business card, notes",
        "  Upload to Claude or ChatGPT",
        "  → \"What's in this image? Describe everything you see.\"",
        "  → \"Read any text in this image and type it out.\"",
        "",
        "This is where AI stops being a toy and starts being a tool"
    ],
    notes="Common reactions: 'Wait, it read my whole 20-page report?' Yes. 'It can read my terrible handwriting?' Mostly, yes. This is the transition from 'AI is fun to play with' to 'AI actually saves me time.'"
)

# ============================================================
# PART 6: OPENCLAW & SETUP
# ============================================================
add_dark_slide("Part 6", "OpenClaw & Setup")

add_content_slide(
    "Meet OpenClaw",
    [
        "One app. Every AI brain. Pay only for what you use.",
        "",
        "Remember the Netflix analogy?",
        "  ChatGPT = Netflix (one company's content)",
        "  OpenClaw = Netflix + Disney+ + HBO in one interface",
        "",
        "✓ Use ANY model — GPT-4o, Claude Sonnet, Gemini — from one place",
        "✓ Runs on your phone, computer, or in Discord",
        "✓ Pay pennies per message through API keys",
        "✓ Switch brains like changing a TV channel"
    ],
    notes="OpenClaw is a personal AI assistant that can use ANY model through one interface. Instead of paying $20/month to each company, you pay pennies per message through API keys."
)

add_table_slide(
    "Real Cost Comparison",
    ["Scenario", "ChatGPT Plus", "Claude Pro", "OpenClaw (API)"],
    [
        ["10 msgs/day, casual use", "$20/mo", "$20/mo", "~$3-5/mo"],
        ["Heavy use + all models", "$20/mo (one company)", "$20/mo (one company)", "~$10-20/mo (ALL)"],
        ["Light use (few times/week)", "$20/mo (overpaying)", "$20/mo (overpaying)", "~$1-2/mo"],
    ],
    notes="Let me show you the model-switching in action. Quick question → cheap model (~$0.0005). Real work → balanced model (~$0.01). Complex analysis → big guns (~$0.05). Match the brain to the task."
)

add_content_slide(
    "Smart Model Switching — Demo",
    [
        "🐇 Quick question → cheap model (~$0.0005):",
        "  /model claude-haiku",
        "  \"What's a good substitute for buttermilk in baking?\"",
        "",
        "🦊 Real work → balanced model (~$0.01):",
        "  /model claude-sonnet",
        "  \"Review this email and make it more professional...\"",
        "",
        "🧠 Complex analysis → max power (~$0.05):",
        "  /model claude-opus",
        "  \"Analyze the financial implications of moving from Toronto to Calgary...\"",
        "",
        "Match the brain to the task. Don't take a taxi to the mailbox."
    ],
    notes="You wouldn't take a taxi to the mailbox, and you wouldn't bike to the airport. Same idea with AI models."
)

# --- NEW: What is a Markdown File? ---
add_content_slide(
    "What is a Markdown File?",
    [
        "A .md file is a plain text file with simple formatting",
        "You can open it in any text editor — it's just text!",
        "",
        "Why should you care?",
        "  • AI tools use markdown files for instructions:",
        "    AGENTS.md, SOUL.md, README.md, TOOLS.md",
        "  • It's how you \"talk\" to AI tools about how to behave",
        "  • It's readable by both humans AND machines",
        "",
        "Think of it as: writing with simple symbols that become formatting"
    ],
    notes="Markdown is the language of AI configuration. When you set up OpenClaw, you'll see files like AGENTS.md — these are markdown files that tell the AI how to behave."
)

add_content_slide(
    "Markdown Basics — The Syntax",
    [
        "# Heading 1          (big title)",
        "## Heading 2         (subtitle)",
        "### Heading 3        (section)",
        "",
        "**bold text**        → bold text",
        "*italic text*        → italic text",
        "",
        "- Bullet point       → • Bullet point",
        "1. Numbered item     → 1. Numbered item",
        "",
        "[link text](url)     → clickable link",
        "`code`               → inline code",
        "```  code block  ``` → multi-line code"
    ],
    notes="You don't need to memorize this. Just know the basics: # for headings, ** for bold, - for bullets, and [text](url) for links. That covers 90% of what you'll see."
)

add_content_slide(
    "Markdown in AI Tools",
    [
        "Common AI configuration files you'll see:",
        "",
        "📄 README.md — Project description (\"what is this?\")",
        "📄 AGENTS.md — Instructions for AI (\"how should you behave?\")",
        "📄 SOUL.md — AI personality/values (\"who are you?\")",
        "📄 TOOLS.md — Available tools (\"what can you use?\")",
        "",
        "These are just text files with .md extension",
        "Edit them in any text editor to customize your AI",
        "",
        "When you edit AGENTS.md, you're literally programming",
        "  your AI assistant — in plain English!"
    ],
    notes="This is the magic of tools like OpenClaw. You don't need to code. You write instructions in plain English in a markdown file, and the AI follows them."
)

# --- NEW: Linux Survival Kit ---
add_content_slide(
    "Linux Survival Kit",
    [
        "What is a terminal?",
        "  A text-based way to talk to your computer",
        "  Instead of clicking icons, you type commands",
        "  It looks scary but you only need ~6 commands",
        "",
        "Why do we need this?",
        "  OpenClaw runs in a terminal (on Windows via WSL)",
        "  Think of it as: the backstage entrance to your computer",
        "",
        "You don't need to be a Linux expert",
        "You just need to not be scared of it"
    ],
    notes="The terminal is just a text-based way to control your computer. It looks like hacker stuff from movies, but really it's just typing simple commands. You'll be fine."
)

add_content_slide(
    "The 6 Commands You Need",
    [
        "pwd          → \"Where am I?\" (print working directory)",
        "ls           → \"What's here?\" (list files in current folder)",
        "cd folder    → \"Go into folder\" (change directory)",
        "cd ..        → \"Go back one level\"",
        "mkdir name   → \"Create a new folder\"",
        "cat file.txt → \"Show me what's in this file\"",
        "",
        "Editing files:",
        "  nano file.txt → Simple text editor (Ctrl+X to exit)",
        "  (vim exists too, but nano is friendlier)",
        "",
        "Pro tip: Tab key auto-completes file/folder names!"
    ],
    notes="That's it. Six commands. pwd to know where you are, ls to see what's there, cd to move around, mkdir to make folders, cat to read files, nano to edit files. Tab auto-completes names so you don't have to type everything."
)

# Setup slides
add_content_slide(
    "Step 1: Get Your API Keys",
    [
        "OpenAI (for GPT models):",
        "  1. Go to platform.openai.com → Sign up",
        "  2. API Keys → Create new secret key → name it \"openclaw\"",
        "  3. COPY IT NOW (you can't see it again!)",
        "  4. Billing → Add $10 credit",
        "",
        "Anthropic (for Claude models):",
        "  1. Go to console.anthropic.com → Sign up",
        "  2. API Keys → Create Key → name it \"openclaw\"",
        "  3. COPY IT NOW → Billing → Add $10 credit",
        "",
        "Google (for Gemini models):",
        "  1. Go to aistudio.google.com → Get API key → Create",
        "  2. COPY IT → generous free tier, may not need credit yet"
    ],
    notes="This feels like a lot of copying and pasting. It is. You do it once and you're done forever. Think of it as setting up your Wi-Fi — annoying for 10 minutes, then you never think about it again."
)

add_content_slide(
    "Step 2: Install OpenClaw — Windows (WSL)",
    [
        "1. Install WSL (open PowerShell as Administrator):",
        "   wsl --install",
        "   → Restart your computer",
        "",
        "2. Open Ubuntu (appears in Start menu)",
        "   → Create username and password",
        "   → Password won't show dots — that's normal!",
        "",
        "3. Update:  sudo apt update && sudo apt upgrade -y",
        "",
        "4. Install Node.js:",
        "   curl -fsSL https://deb.nodesource.com/setup_22.x | sudo -E bash -",
        "   sudo apt install -y nodejs",
        "",
        "5. Install OpenClaw:  npm install -g @anthropic/openclaw",
        "6. Setup:  openclaw setup  (paste your API keys)"
    ],
    notes="Windows users need WSL first. WSL = Windows Subsystem for Linux — a little Linux computer inside your Windows computer. This is the most annoying part. Once it's done, everything else is easy."
)

add_content_slide(
    "Step 2: Install OpenClaw — macOS",
    [
        "1. Open Terminal (Cmd+Space → type \"Terminal\" → Enter)",
        "",
        "2. Install Homebrew (if you don't have it):",
        "   /bin/bash -c \"$(curl -fsSL https://raw.githubusercontent.com/",
        "   Homebrew/install/HEAD/install.sh)\"",
        "",
        "3. Install Node.js:  brew install node",
        "",
        "4. Install OpenClaw:  npm install -g @anthropic/openclaw",
        "",
        "5. Setup:  openclaw setup  (paste your API keys)"
    ],
    notes="macOS is simpler. Terminal is already there. Homebrew is the package manager. Then Node.js, then OpenClaw."
)

add_content_slide(
    "Step 3: Verify It Works",
    [
        "In your terminal, type:",
        "",
        "  openclaw chat \"Hello! Tell me a fun fact.\"",
        "",
        "If you see a response — you're done! 🎉",
        "",
        "Common issues:",
        "  \"command not found\" → Node didn't install. Run: node --version",
        "  \"invalid API key\" → Copied wrong. Create a new one.",
        "  \"insufficient funds\" → Add credit to your API account.",
        "",
        "You now have every major AI brain available from one tool."
    ],
    notes="Hands up if you got a response! If it didn't work, don't panic. I'll stay after the session to help anyone who's stuck."
)

# Cheat Sheet / Q&A
add_content_slide(
    "What to Remember from Today",
    [
        "1. There are 3 major AIs — ChatGPT, Claude, Gemini",
        "2. The App ≠ The Brain — apps deliver, models think",
        "3. Match the model to the task (fast/cheap vs smart/expensive)",
        "4. AI hallucinates — always verify important facts",
        "5. Privacy: treat AI like a coworker you don't fully trust",
        "6. API keys = pay-per-use, cheaper than subscriptions",
        "7. OpenClaw = one tool for all AI brains",
        "",
        "Start with Claude Sonnet for most things",
        "GPT-4o when you need images | Gemini for current info"
    ],
    notes="You walked in knowing ChatGPT. You're walking out knowing how to use ANY AI, switch between them, and pay a fraction of what a subscription costs. That's real power."
)

add_dark_slide(
    "Questions?",
    "Nothing is too basic — if you're wondering it, three other people are too."
)

# Save
output_path = "/home/lj_wsl/cut-the-crap/general/cut-the-crap-general.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
